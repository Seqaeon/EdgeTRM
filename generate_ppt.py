import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper: Set dark background
    def add_background(slide):
        rect = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE = 1
            0, 0, prs.slide_width, prs.slide_height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
        rect.line.fill.background() # No border
        return rect

    # Helper: Add slide title
    def add_slide_header(slide, title, category="TECHNICAL BRIEF"):
        # Category label (small, top)
        catBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(11.83), Inches(0.4))
        cat_tf = catBox.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = "Arial"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = RGBColor(6, 182, 212) # Cyan 500

        # Slide Title
        titleBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.83), Inches(0.8))
        title_tf = titleBox.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.name = "Arial"
        title_p.font.size = Pt(26)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)

    # Helper: Add card and text
    def add_card(slide, left, top, width, height, title, lines):
        # Base Card Shape
        card = slide.shapes.add_shape(5, left, top, width, height) # MSO_SHAPE.ROUNDED_RECTANGLE = 5
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
        card.line.color.rgb = RGBColor(51, 65, 85) # Slate 700
        card.line.width = Pt(1.5)

        # Text Frame Box
        pad = Inches(0.25)
        txBox = slide.shapes.add_textbox(left + pad, top + pad, width - (pad * 2), height - (pad * 2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        # Title of the Card
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(6, 182, 212) # Cyan 500
        p.space_after = Pt(12)

        for line in lines:
            p = tf.add_paragraph()
            p.font.name = "Arial"
            p.space_after = Pt(6)
            
            if line.startswith("• "):
                p.text = line[2:]
                p.level = 0
                p.font.size = Pt(12.5)
                p.font.color.rgb = RGBColor(255, 255, 255)
            elif line.startswith("  - "):
                p.text = line[4:]
                p.level = 1
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
            elif line.startswith("💡 ") or line.startswith("⚠️ ") or line.startswith("✅ ") or line.startswith("❌ "):
                p.text = line
                p.font.size = Pt(12)
                p.font.italic = True
                p.font.color.rgb = RGBColor(6, 182, 212) # Cyan
                p.space_before = Pt(6)
            else:
                p.text = line
                p.font.size = Pt(12.5)
                p.font.color.rgb = RGBColor(255, 255, 255)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide_1 = prs.slides.add_slide(blank_layout)
    add_background(slide_1)

    # Decorative vertical accent bar
    bar = slide_1.shapes.add_shape(1, Inches(0.75), Inches(2.2), Inches(0.12), Inches(3.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(6, 182, 212)
    bar.line.fill.background()

    # Text box
    tx_title = slide_1.shapes.add_textbox(Inches(1.15), Inches(2.0), Inches(11.0), Inches(4.0))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0

    p_main = tf_title.paragraphs[0]
    p_main.text = "EdgeTRM: Compressing 68.89% ARC-AGI Models for Edge Hardware"
    p_main.font.name = "Arial"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = RGBColor(255, 255, 255)
    p_main.space_after = Pt(14)

    p_sub = tf_title.add_paragraph()
    p_sub.text = "A Technical Briefing on Model Footprint, Compute Efficiency, and Reasoning Trade-Offs"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(148, 163, 184)
    p_sub.space_after = Pt(28)

    p_footer = tf_title.add_paragraph()
    p_footer.text = "EdgeTRM Compression & Deployment Analysis Team"
    p_footer.font.name = "Arial"
    p_footer.font.size = Pt(13)
    p_footer.font.bold = True
    p_footer.font.color.rgb = RGBColor(6, 182, 212)

    # ----------------------------------------------------
    # SLIDE 2: Memory Bottlenecks & SRAM Coexistence
    # ----------------------------------------------------
    slide_2 = prs.slides.add_slide(blank_layout)
    add_background(slide_2)
    add_slide_header(slide_2, "01. Memory Bottlenecks & SRAM Coexistence")

    card_y = Inches(1.6)
    card_h = Inches(5.1)
    card_w = Inches(3.6)
    gap = Inches(0.5)

    add_card(
        slide_2, Inches(0.75), card_y, card_w, card_h,
        "The Memory Bottleneck",
        [
            "• Baseline FP32 footprint is **88 MB**.",
            "• **Puzzle Embeddings (61 MB):**",
            "  - Dominates 70% of total footprint.",
            "  - Completely blocks edge deployment.",
            "• **Backbone (26.7 MB):**",
            "  - Active neural weights.",
            "💡 Puzzle embeddings are the critical barrier to fit inside SRAM limit."
        ]
    )

    add_card(
        slide_2, Inches(0.75) + card_w + gap, card_y, card_w, card_h,
        "Quantization Frontier",
        [
            "• **INT8 (bitsandbytes):**",
            "  - Shrinks backbone to **6.7 MB**.",
            "  - Accuracy actually improves to 20.03% (vs 19.20% FP32 baseline).",
            "• **Calibrated INT4:**",
            "  - Shrinks backbone to **3.3 MB**.",
            "  - Recovers baseline carry fidelity and achieves 19.02% cell accuracy.",
            "❌ Naive INT4 collapses performance entirely."
        ]
    )

    add_card(
        slide_2, Inches(0.75) + (card_w + gap) * 2, card_y, card_w, card_h,
        "SRAM Coexistence",
        [
            "• **Single-Puzzle Loading:**",
            "  - Keep the massive 61 MB embedding table on external flash storage.",
            "  - Load only the active **2 KB row** for the current puzzle into SRAM.",
            "  - Achieves zero information loss.",
            "✅ **Fits 8 MB SRAM:**",
            "  - Fused INT8 backbone + active row uses only **6.59 MB active SRAM**."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 3: Reasoning Fidelity & 1-Cycle Compute Optimization
    # ----------------------------------------------------
    slide_3 = prs.slides.add_slide(blank_layout)
    add_background(slide_3)
    add_slide_header(slide_3, "02. Carry-State Fidelity & 1-Cycle Compute Optimization")

    card_w_2 = Inches(5.6)
    gap_2 = Inches(0.63)

    add_card(
        slide_3, Inches(0.75), card_y, card_w_2, card_h,
        "Carry-State Reasoning Health",
        [
            "Using carry-state cosine similarity to diagnose reasoning health across recursive depth steps:",
            "• **INT8 & FP32 (0.377 similarity):**",
            "  - Model actively refines carry states per step; reasoning is intact.",
            "• **Calibrated INT4 (0.975 similarity to FP32):**",
            "  - Tracks FP32 reasoning trajectory faithfully.",
            "• **Naive INT4 (0.722 similarity):**",
            "  - States freeze/collapse. Reasoning fails entirely.",
            "💡 INT8 error accumulation explains why multi-step inference degrades while 1-step stays high."
        ]
    )

    add_card(
        slide_3, Inches(0.75) + card_w_2 + gap_2, card_y, card_w_2, card_h,
        "1-Cycle Compute Breakthrough",
        [
            "• **16× Compute Reduction:**",
            "  - Cutting recursive depth from 16 to 1 step saves 2,812 GFLOPs per puzzle (3,000 → 187 GFLOPs).",
            "• **Superior Accuracy:**",
            "  - INT8 cell accuracy at 1 H-cycle is actually **higher** (19.24%) than at 16 H-cycles (19.00%).",
            "  - FP16 & FP32 maintain strong performance at 1 cycle.",
            "✅ **Production Ready:**",
            "  - Instant 16× cost reduction in production with zero retraining."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 4: Deployment Feasibility & Strategic Roadmap
    # ----------------------------------------------------
    slide_4 = prs.slides.add_slide(blank_layout)
    add_background(slide_4)
    add_slide_header(slide_4, "03. Deployment Feasibility & Strategic Roadmap")

    add_card(
        slide_4, Inches(0.75), card_y, card_w_2, card_h,
        "The MCU Compute Barrier",
        [
            "SRAM footprints fit 8 MB, but the O(n²) attention computation over 900 tokens blocks microcontrollers (MCUs):",
            "• **Cortex-A55 (Mobile SoC):** ~93s per puzzle.",
            "• **Cortex-M55 (Edge MCU):** ~622s per puzzle.",
            "• **ESP32-S3 (Edge MCU):** ~1,554s per puzzle.",
            "⚠️ MCU-class bare-metal execution is borderline impossible without sequence length reduction."
        ]
    )

    add_card(
        slide_4, Inches(0.75) + card_w_2 + gap_2, card_y, card_w_2, card_h,
        "Strategic Recommendations",
        [
            "1. **Shift Target to Mobile SoCs:**",
            "   - Target Snapdragon or Raspberry Pi (4-8 GB RAM makes memory trivial; CPU compute is viable).",
            "2. **Lock 1-Cycle Inference:**",
            "   - Instantly save 16× compute in production immediately using H_cycles=1.",
            "3. **Native INT8 Transition:**",
            "   - Convert from bitsandbytes to PyTorch-native `torch.ao.quantization` for CPU tracing.",
            "4. **Reduce Context Length:**",
            "   - Crop/stride to 450 tokens to save 4-8× compute."
        ]
    )

    output_path = "/home/seqaeon/Downloads/EdgeTRM/EdgeTRM_Compression_Brief.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
